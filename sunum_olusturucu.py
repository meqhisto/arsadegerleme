"""
Arsa Yatırım Danışmanlığı - Sunum Oluşturucu Modülü
Bu modül, arsa analiz sonuçlarından Word ve PowerPoint sunumları oluşturmak için kullanılır.
"""

import os
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

class SunumOlusturucu:
    def __init__(self, templates_dir):
        """
        Sunum oluşturucu sınıfını başlatır.
        
        Args:
            templates_dir (str): Şablon dosyalarının bulunduğu dizin
        """
        self.templates_dir = templates_dir
        os.makedirs(templates_dir, exist_ok=True)
        
        # Şablon dosya yolları
        self.word_template = os.path.join(templates_dir, "word_template.docx")
        self.ppt_template = os.path.join(templates_dir, "ppt_template.pptx")
        
        # Şablonlar yoksa oluştur
        self._create_templates_if_not_exist()
    
    def _create_templates_if_not_exist(self):
        """
        Şablon dosyaları yoksa oluşturur.
        """
        # Word şablonu oluştur
        if not os.path.exists(self.word_template):
            self._create_word_template()
        
        # PowerPoint şablonu oluştur
        if not os.path.exists(self.ppt_template):
            self._create_ppt_template()
    
    def _create_word_template(self):
        """
        Word şablonu oluşturur.
        """
        doc = Document()
        
        # Sayfa yapısını ayarla
        sections = doc.sections
        for section in sections:
            section.page_height = Cm(29.7)
            section.page_width = Cm(21)
            section.left_margin = Cm(2.5)
            section.right_margin = Cm(2.5)
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
        
        # Başlık ekle
        title = doc.add_heading("ARSA YATIRIM ANALİZ RAPORU", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Alt başlıklar ve yer tutucular ekle
        doc.add_heading("1. ARSA BİLGİLERİ", level=1)
        doc.add_paragraph("Konum: {konum}")
        doc.add_paragraph("Metrekare: {metrekare} m²")
        doc.add_paragraph("İmar Durumu: {imar_durumu}")
        doc.add_paragraph("Fiyat: {fiyat} TL")
        doc.add_paragraph("Metrekare Fiyatı: {metrekare_fiyat} TL/m²")
        doc.add_paragraph("Bölge Ortalama Fiyatı: {bolge_fiyat} TL/m²")
        
        doc.add_heading("2. ANALİZ SONUÇLARI", level=1)
        doc.add_paragraph("Bölge Karşılaştırması: {bolge_karsilastirma}%")
        doc.add_paragraph("Potansiyel Getiri: %{potansiyel_getiri}")
        doc.add_paragraph("Tavsiye Edilen Yatırım Süresi: {yatirim_suresi} yıl")
        
        doc.add_heading("3. YATIRIM DEĞERLENDİRMESİ", level=1)
        doc.add_paragraph("{temel_ozet}")
        doc.add_paragraph("{yatirim_ozet}")
        
        doc.add_heading("4. TAVSİYELER", level=1)
        doc.add_paragraph("{tavsiyeler}")
        
        # Şablonu kaydet
        doc.save(self.word_template)
    
    def _create_ppt_template(self):
        """
        PowerPoint şablonu oluşturur.
        """
        prs = Presentation()
        
        # Kapak slaytı
        slide_layout = prs.slide_layouts[0]  # Başlık slaytı
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "ARSA YATIRIM ANALİZ RAPORU"
        subtitle.text = "{konum}"
        
        # Arsa bilgileri slaytı
        slide_layout = prs.slide_layouts[1]  # Başlık ve içerik slaytı
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "ARSA BİLGİLERİ"
        content.text = "• Konum: {konum}\n"
        content.text += "• Metrekare: {metrekare} m²\n"
        content.text += "• İmar Durumu: {imar_durumu}\n"
        content.text += "• Fiyat: {fiyat} TL\n"
        content.text += "• Metrekare Fiyatı: {metrekare_fiyat} TL/m²\n"
        content.text += "• Bölge Ortalama Fiyatı: {bolge_fiyat} TL/m²"
        
        # Analiz sonuçları slaytı
        slide_layout = prs.slide_layouts[1]  # Başlık ve içerik slaytı
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "ANALİZ SONUÇLARI"
        content.text = "• Bölge Karşılaştırması: {bolge_karsilastirma}%\n"
        content.text += "• Potansiyel Getiri: %{potansiyel_getiri}\n"
        content.text += "• Tavsiye Edilen Yatırım Süresi: {yatirim_suresi} yıl"
        
        # Yatırım değerlendirmesi slaytı
        slide_layout = prs.slide_layouts[1]  # Başlık ve içerik slaytı
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "YATIRIM DEĞERLENDİRMESİ"
        content.text = "{temel_ozet}\n\n"
        content.text += "{yatirim_ozet}"
        
        # Tavsiyeler slaytı
        slide_layout = prs.slide_layouts[1]  # Başlık ve içerik slaytı
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "TAVSİYELER"
        content.text = "{tavsiyeler}"
        
        # Şablonu kaydet
        prs.save(self.ppt_template)
    
    def olustur_word(self, arsa_data, output_path):
        """
        Word belgesi oluşturur.
        
        Args:
            arsa_data (dict): Arsa verileri
            output_path (str): Çıktı dosyasının yolu
            
        Returns:
            str: Oluşturulan dosyanın yolu
        """
        # Şablonu yükle
        doc = Document(self.word_template)
        
        # Değerleri formatla
        konum = f"{arsa_data['konum']['mahalle']}, {arsa_data['konum']['ilce']}, {arsa_data['konum']['il']}"
        metrekare = f"{arsa_data['metrekare']:,.2f}"
        imar_durumu = arsa_data['imar_durumu'].capitalize()
        fiyat = f"{arsa_data['fiyat']:,.2f}"
        metrekare_fiyat = f"{arsa_data['metrekare_fiyat']:,.2f}"
        bolge_fiyat = f"{arsa_data['bolge_fiyat']:,.2f}"
        
        bolge_karsilastirma = f"{arsa_data['bolge_karsilastirma']:.1f}"
        if arsa_data['bolge_karsilastirma'] < 0:
            bolge_karsilastirma = f"-{abs(arsa_data['bolge_karsilastirma']):.1f}"
        
        potansiyel_getiri = f"{arsa_data['potansiyel_getiri']:.1f}"
        yatirim_suresi = str(arsa_data['yatirim_suresi'])
        
        # Özet metinleri
        temel_ozet = arsa_data.get('ozet', {}).get('temel_ozet', '')
        yatirim_ozet = arsa_data.get('ozet', {}).get('yatirim_ozet', '')
        tavsiyeler = arsa_data.get('ozet', {}).get('tavsiyeler', '')
        
        # Belgedeki yer tutucuları değiştir
        for paragraph in doc.paragraphs:
            paragraph.text = paragraph.text.replace("{konum}", konum)
            paragraph.text = paragraph.text.replace("{metrekare}", metrekare)
            paragraph.text = paragraph.text.replace("{imar_durumu}", imar_durumu)
            paragraph.text = paragraph.text.replace("{fiyat}", fiyat)
            paragraph.text = paragraph.text.replace("{metrekare_fiyat}", metrekare_fiyat)
            paragraph.text = paragraph.text.replace("{bolge_fiyat}", bolge_fiyat)
            paragraph.text = paragraph.text.replace("{bolge_karsilastirma}", bolge_karsilastirma)
            paragraph.text = paragraph.text.replace("{potansiyel_getiri}", potansiyel_getiri)
            paragraph.text = paragraph.text.replace("{yatirim_suresi}", yatirim_suresi)
            paragraph.text = paragraph.text.replace("{temel_ozet}", temel_ozet)
            paragraph.text = paragraph.text.replace("{yatirim_ozet}", yatirim_ozet)
            paragraph.text = paragraph.text.replace("{tavsiyeler}", tavsiyeler)
        
        # Belgeyi kaydet
        doc.save(output_path)
        
        return output_path
    
    def olustur_powerpoint(self, arsa_data, output_path):
        """
        PowerPoint sunumu oluşturur.
        
        Args:
            arsa_data (dict): Arsa verileri
            output_path (str): Çıktı dosyasının yolu
            
        Returns:
            str: Oluşturulan dosyanın yolu
        """
        # Şablonu yükle
        prs = Presentation(self.ppt_template)
        
        # Değerleri formatla
        konum = f"{arsa_data['konum']['mahalle']}, {arsa_data['konum']['ilce']}, {arsa_data['konum']['il']}"
        metrekare = f"{arsa_data['metrekare']:,.2f}"
        imar_durumu = arsa_data['imar_durumu'].capitalize()
        fiyat = f"{arsa_data['fiyat']:,.2f}"
        metrekare_fiyat = f"{arsa_data['metrekare_fiyat']:,.2f}"
        bolge_fiyat = f"{arsa_data['bolge_fiyat']:,.2f}"
        
        bolge_karsilastirma = f"{arsa_data['bolge_karsilastirma']:.1f}"
        if arsa_data['bolge_karsilastirma'] < 0:
            bolge_karsilastirma = f"-{abs(arsa_data['bolge_karsilastirma']):.1f}"
        
        potansiyel_getiri = f"{arsa_data['potansiyel_getiri']:.1f}"
        yatirim_suresi = str(arsa_data['yatirim_suresi'])
        
        # Özet metinleri
        temel_ozet = arsa_data.get('ozet', {}).get('temel_ozet', '')
        yatirim_ozet = arsa_data.get('ozet', {}).get('yatirim_ozet', '')
        tavsiyeler = arsa_data.get('ozet', {}).get('tavsiyeler', '')
        
        # Slaytlardaki yer tutucuları değiştir
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape.text = shape.text.replace("{konum}", konum)
                    shape.text = shape.text.replace("{metrekare}", metrekare)
                    shape.text = shape.text.replace("{imar_durumu}", imar_durumu)
                    shape.text = shape.text.replace("{fiyat}", fiyat)
                    shape.text = shape.text.replace("{metrekare_fiyat}", metrekare_fiyat)
                    shape.text = shape.text.replace("{bolge_fiyat}", bolge_fiyat)
                    shape.text = shape.text.replace("{bolge_karsilastirma}", bolge_karsilastirma)
                    shape.text = shape.text.replace("{potansiyel_getiri}", potansiyel_getiri)
                    shape.text = shape.text.replace("{yatirim_suresi}", yatirim_suresi)
                    shape.text = shape.text.replace("{temel_ozet}", temel_ozet)
                    shape.text = shape.text.replace("{yatirim_ozet}", yatirim_ozet)
                    shape.text = shape.text.replace("{tavsiyeler}", tavsiyeler)
        
        # Sunumu kaydet
        prs.save(output_path)
        
        return output_path
